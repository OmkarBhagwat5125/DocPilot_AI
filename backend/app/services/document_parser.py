import os
import re
import logging
import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
from pathlib import Path

logger = logging.getLogger(__name__)

class DocumentParser:
    def parse(self, file_path: str, filename: str) -> list[dict]:
        path = Path(file_path)
        ext = path.suffix.lower()
        
        logger.info(f"Parsing document {filename} of type {ext}")
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        try:
            if ext == ".pdf":
                return self._parse_pdf(file_path, filename)
            elif ext == ".docx":
                return self._parse_docx(file_path, filename)
            elif ext == ".pptx":
                return self._parse_pptx(file_path, filename)
            elif ext == ".csv":
                return self._parse_csv(file_path, filename)
            elif ext == ".txt":
                return self._parse_txt(file_path, filename)
            else:
                raise ValueError(f"Unsupported file extension: {ext}")
        except Exception as e:
            logger.error(f"Error parsing file {filename}: {str(e)}")
            raise e

    def _clean_text(self, text: str) -> str:
        if not text:
            return ""
        # Remove duplicate whitespace/newlines
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _parse_pdf(self, file_path: str, filename: str) -> list[dict]:
        results = []
        doc = fitz.open(file_path)
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            text = page.get_text()
            cleaned = self._clean_text(text)
            if cleaned:
                results.append({
                    "text": cleaned,
                    "page_number": page_idx + 1,
                    "source": filename
                })
        doc.close()
        return results

    def _parse_docx(self, file_path: str, filename: str) -> list[dict]:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
                
        # Group into a single document object for now (or split by custom section separator if needed)
        cleaned = self._clean_text("\n".join(full_text))
        if cleaned:
            return [{
                "text": cleaned,
                "page_number": 1,
                "source": filename
            }]
        return []

    def _parse_pptx(self, file_path: str, filename: str) -> list[dict]:
        prs = pptx.Presentation(file_path)
        results = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            cleaned = self._clean_text(" ".join(slide_text))
            if cleaned:
                results.append({
                    "text": cleaned,
                    "page_number": slide_idx + 1,
                    "source": filename
                })
        return results

    def _parse_csv(self, file_path: str, filename: str) -> list[dict]:
        df = pd.read_csv(file_path)
        rows_text = []
        for idx, row in df.iterrows():
            row_str = ", ".join([f"{col}: {val}" for col, val in row.items()])
            rows_text.append(row_str)
            
        cleaned = self._clean_text("\n".join(rows_text))
        if cleaned:
            return [{
                "text": cleaned,
                "page_number": 1,
                "source": filename
            }]
        return []

    def _parse_txt(self, file_path: str, filename: str) -> list[dict]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        cleaned = self._clean_text(text)
        if cleaned:
            return [{
                "text": cleaned,
                "page_number": 1,
                "source": filename
            }]
        return []
